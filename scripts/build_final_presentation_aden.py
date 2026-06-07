#!/usr/bin/env python3
"""Generate Final_Presentation_Aden.pptx for Week 10 capstone."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "Final_Presentation_Aden.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x0D, 0x94, 0x88)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)

FONT_NAME = "Calibri"
CONTENT_TITLE_SIZE = 28
CONTENT_BODY_SIZE = 23


def apply_font(paragraph, *, size, bold=True, color=None):
    paragraph.font.name = FONT_NAME
    paragraph.font.bold = bold
    paragraph.font.size = Pt(size)
    if color is not None:
        paragraph.font.color.rgb = color


def style_content_title(slide):
    for p in slide.shapes.title.text_frame.paragraphs:
        apply_font(p, size=CONTENT_TITLE_SIZE, color=NAVY)


def set_title(slide, text, subtitle=None):
    slide.shapes.title.text = text
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = NAVY
        elif shape.placeholder_format.idx == 1:
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = MUTED


def add_bullets(slide, items):
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(items):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        apply_font(p, size=CONTENT_BODY_SIZE, color=SLATE)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]

    slides_data = [
        {
            "layout": layout_title,
            "title": "Job Application Tracker",
            "subtitle": "High Flyers — Week 10 Final Presentation\nAden Abdulahi · DB & Security\nLive: https://18.118.119.37 · Repo: course-506-high-flyers/job-tracker",
            "notes": "~30 sec. Introduce team (Boma, Darrell, Aden), project name, live URL.",
        },
        {
            "title": "What We Built (~1 min)",
            "bullets": [
                "Web app for job seekers tracking applications in one place",
                "GitHub OAuth login, full CRUD on applications, stage tracking",
                "Company enrichment via Clearbit + API Ninjas (cached per company)",
                "Deployed on AWS EC2: nginx → gunicorn → Flask → Postgres",
                "Who it's for: anyone applying to multiple roles who needs pipeline visibility",
            ],
            "notes": "Elevator pitch only — one minute. Don't demo yet.",
        },
        {
            "title": "Architecture (~1.5 min)",
            "bullets": [
                "Browser → nginx (:443 TLS, :80 redirect) → gunicorn (:8000) → Flask app",
                "Flask ↔ Postgres 16 (SQLModel ORM, named volume pgdata)",
                "Three Docker containers: nginx, app, db — docker compose up -d",
                "Contracts: CONTRACTS.md is the authoritative spec all slices follow",
                "Edge owns TLS, rate limits, attack-path 404s; app owns auth + business logic",
            ],
            "notes": "Point at diagram on slide. Trace one request: HTTPS hits nginx, proxy_pass to app:8000, Flask queries db on internal network only.",
        },
        {
            "title": "Architecture Diagram",
            "bullets": [
                "[Internet] → [nginx :443] → [gunicorn :8000 / Flask app]",
                "                                    ↓",
                "                              [Postgres :5432 internal]",
                "                                    ↓",
                "                         [Clearbit / API Ninjas — insight only]",
                "Key contracts: X-Forwarded-Proto + ProxyFix · DATABASE_URL · OAuth secrets via .env",
            ],
            "notes": "Walk left-to-right. Emphasize db has no published host port — network trust boundary.",
        },
        {
            "title": "My Slice: DB & Security (~1.5 min)",
            "bullets": [
                "nginx.conf — TLS, 4 rate-limit zones, attack-path 404s before Flask",
                "Proved empirically: sidecar test, 16 probes, Flask request-count delta = 0",
                "models.py — unique constraints, CASCADE deletes, BOLA 404 (not 403)",
                "tests/test_attack_paths.py + attack_paths.json — 33 automated regression tests",
                "Secrets hygiene: .gitignore for .env.save*, strict env loading, dev cert script",
            ],
            "notes": "This is your strongest Week 8 evidence. Mention instructor praised the empirical nginx proof.",
        },
        {
            "title": "Live Demo Script (~2 min)",
            "bullets": [
                "1. Open https://18.118.119.37 — accept self-signed cert if prompted",
                "2. Sign in with GitHub OAuth → lands on My Applications",
                "3. Add application (company, role, date) → view detail page",
                "4. Trigger company insight — external API enrichment",
                "5. Security quick-check: curl /.env → 404 from nginx, never hits Flask",
                "If something breaks: show docker compose ps + logs — honest recovery",
            ],
            "notes": "Demo is centerpiece. Have terminal ready for attack-path curl as backup if UI hiccups.",
        },
        {
            "title": "What Worked · What Didn't (~1 min)",
            "bullets": [
                ("Solid:", 0),
                ("CONTRACTS.md kept three slices integrating without schema churn", 1),
                ("nginx attack-path filtering — measured, not assumed", 1),
                ("GitHub OAuth + ProxyFix after HTTPS migration", 1),
                ("Held together with tape:", 0),
                ("Single EC2 + single pgdata volume — no backup runbook yet", 1),
                ("Mutable Docker image tags — not digest-pinned", 1),
                ("Would redo: CI deploy workflow earlier; pg_dump cron to S3", 1),
            ],
            "notes": "Honesty scores points. Don't oversell.",
        },
        {
            "title": "Team Split & Git Workflow (~45 sec)",
            "bullets": [
                "Boma — templates, Bootstrap, security headers, coordination",
                "Darrell — Flask routes, gunicorn, ProxyFix, external API integration",
                "Aden — schema, auth hardening, nginx edge, attack-path tests, deploy",
                "Flow: personal branches → hardening → main via PR (#12–#14)",
                "Branch protection: 1 approving review required on main",
            ],
            "notes": "Quick handoff — team knows this. Point at team_evidence.md PR table if asked.",
        },
        {
            "title": "AI Usage Reflection (~1 min)",
            "bullets": [
                ("Where AI helped:", 0),
                ("Scaffolding nginx location blocks and pytest parametrize tests", 1),
                ("LLM surface probe surfaced gaps (mutable tags, backup SPOF)", 1),
                ("Where AI didn't:", 0),
                ("X-Forwarded-Proto — said Flask picks it up automatically; needed ProxyFix", 1),
                ("OAuth callback URL mismatch after HTTPS deploy — had to verify manually", 1),
                ("Discipline: one-pass LLM probes, then verify every claim against real files", 1),
            ],
            "notes": "Course is about engineering WITH AI — show triage, not blind trust.",
        },
        {
            "title": "Thank You · Q&A",
            "bullets": [
                "Live app: https://18.118.119.37",
                "GitHub: github.com/course-506-high-flyers/job-tracker",
                "Questions?",
            ],
            "notes": "~30 sec close. Total deck target: 8–9 minutes.",
        },
    ]

    for i, data in enumerate(slides_data):
        if data.get("layout") == layout_title:
            slide = prs.slides.add_slide(layout_title)
            set_title(slide, data["title"], data.get("subtitle"))
        else:
            slide = prs.slides.add_slide(layout_content)
            slide.shapes.title.text = data["title"]
            style_content_title(slide)
            add_bullets(slide, data["bullets"])
        add_notes(slide, data["notes"])

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
